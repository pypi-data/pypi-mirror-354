from os.path import isfile

import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader


def read_file(file_name: str):
    """
    Reads the information from a file
    :param file_name:
    str: Location and name of file
    :return:
    str: The extracted text
    """
    # Check that it is a file
    if not isfile(file_name):
        raise ValueError("'{0}' is not a file".format(file_name))

    # Initialize text from file
    text = []

    # Handle PDFs
    if file_name.endswith(".pdf"):
        for page in PdfReader(file_name).pages:
            text.append(page.extract_text())
    # Handle CSVs
    elif file_name.endswith(".csv"):
        text.append(pd.read_csv(file_name).to_string())
    # Handle pptx
    elif file_name.endswith(".pptx"):
        for slide in Presentation(file_name).slides:
            if not slide.has_text_frame:
                continue
            for paragraph in slide.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
    # Handle docx
    elif file_name.endswith(".docx"):
        for paragraph in Document(file_name).paragraphs:
            for run in paragraph.runs:
                text.append(run.text)
    # Read file and extract text
    else:
        try:
            with open(file_name, "r") as input_file:
                text.append(input_file.read())
        except Exception:
            raise ValueError("File could not be read: {0}".format(file_name))

    print("Read {0}".format(file_name))
    return "\n".join(text)


def reformat_logs(logs: list, new_query_text: str = "-" * 15):
    """
    Reformats the logs from ollama's format to a more user-friendly one for saving to text

    :param logs:
    list: A list of logs to be reformatted. Only reformats the ollama log format, i.e. a list of dicts with `role` and `content` keys.
    :param new_query_text:
    str: Text used to break up the different queries and responses
    :return:
    list: The reformatted logs
    """
    # Initialize logs
    logs_formatted = []

    # Run through each log
    for log in logs:
        # Ensure that the log is in the correct format
        if not isinstance(log, dict) or "role" not in log or "content" not in log:
            raise TypeError("Cannot reformat logs")
        # Get log from user format
        elif log["role"] == "user":
            logs_formatted.append(
                "{0}\nQuery: {1}\n{0}".format(new_query_text, log["content"])
            )
        # Get log from assistant format
        else:
            logs_formatted.append(log["content"])

    return logs_formatted


def unformat_logs(logs: str, new_query_text: str = "-" * 15):
    """
    Unformats the logs from a string into a list of dicts with `role` and `content` keys. This can then be read in with ollama.

    :param logs:
    str: The logs from a text file
    :param new_query_text:
    str: Text used to break up the different queries and responses
    :return:
    list: The unformatted logs in the ollama log format, i.e. a list of dicts with `role` and `content` keys.
    """
    # Initialize logs
    logs_unformatted = []

    # Run through each query/response
    for content in logs.split(new_query_text):
        # Strip content
        content = content.strip()

        # Ensure there is actually a query/response
        if len(content) > 0:
            # Get role
            if content.startswith("Query: "):
                role = "user"
                content = content.replace("Query: ", "")
            else:
                role = "assistant"

            # Add query/response
            logs_unformatted.append({"role": role, "content": content})

    return logs_unformatted
