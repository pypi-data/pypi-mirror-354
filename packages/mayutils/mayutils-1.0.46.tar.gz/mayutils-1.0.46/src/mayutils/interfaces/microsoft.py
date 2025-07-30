from pathlib import Path
from typing import Optional, Self
from pptx import Presentation as Init
from pptx.shapes.autoshape import Shape
from pptx.slide import SlideLayout, Slides, Slide
from pptx.util import Pt, Length as BaseLength
from pptx.dml.color import RGBColor
from mayutils.objects.colours import Colour
import six
from copy import deepcopy


class Length(BaseLength):
    @classmethod
    def from_float(
        cls,
        value: float,
    ) -> Self:
        return cls(value)  # type: ignore


class Presentation:
    def __init__(
        self,
        template: Path | str,
    ) -> None:
        self.template = Path(template)
        if not self.template.exists():
            raise FileNotFoundError(f"Template file {self.template} does not exist.")
        elif self.template.is_dir():
            raise ValueError(
                f"Template file {self.template} is a directory, not a file."
            )
        elif not self.template.is_file():
            raise ValueError(f"Template file {self.template} is not a valid file.")
        elif self.template.suffix.lower() not in [".pptx", ".ppt"]:
            raise ValueError(
                f"Template file {self.template} is not a valid PowerPoint file."
            )
        elif self.template.suffix.lower() == ".ppt":
            raise ValueError(
                f"Template file {self.template} is a legacy PowerPoint file (.ppt). "
                "Please convert it to .pptx format."
            )

        self.internal = Init(pptx=str(self.template))

        self.blank_layout = self.internal.slide_layouts[
            len(self.internal.slide_layouts) - 1
        ]

    @property
    def height(
        self,
    ) -> Length | None:
        return (
            Length(emu=self.internal.slide_height)
            if self.internal.slide_height is not None
            else None
        )

    @height.setter
    def height(
        self,
        value: Length,
    ) -> None:
        self.internal.slide_height = value

    @property
    def width(
        self,
    ) -> Length | None:
        return (
            Length(emu=self.internal.slide_width)
            if self.internal.slide_width is not None
            else None
        )

    @width.setter
    def width(
        self,
        value: Length,
    ) -> None:
        self.internal.slide_width = value

    @property
    def slides(
        self,
    ) -> Slides:
        return self.internal.slides

    def slide(
        self,
        slide_number: int,
    ) -> Slide:
        if slide_number < 1 or slide_number > len(self.slides):
            raise IndexError(
                f"Slide number {slide_number} is out of range. Presentation has {len(self.slides)} slides."
            )

        return self.slides[slide_number - 1]

    def new_slide(
        self,
        layout: Optional[SlideLayout] = None,
    ) -> Self:
        self.slides.add_slide(
            slide_layout=layout if layout is not None else self.blank_layout
        )

        return self

    def delete_slide(
        self,
        slide_number: int,
    ) -> Self:
        raise NotImplementedError("Deleting slides is not implemented yet.")
        return self

    def copy_slide(
        self,
        slide_number: int,
    ) -> Self:
        raise NotImplementedError("Copying slides is not implemented yet.")
        slide_idx = slide_number - 1
        template_slide = self.slides[slide_idx]

        try:
            blank_slide_layout = self.internal.slide_layouts[12]
        except IndexError:
            blank_slide_layout = self.internal.slide_layouts[
                len(self.internal.slide_layouts) - 1
            ]

        copied_slide = self.slides.add_slide(slide_layout=blank_slide_layout)

        for shape in template_slide.shapes:
            element = shape.element
            new_element = deepcopy(element)
            copied_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

        for _, value in six.iteritems(template_slide.part.rels):
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(
                    value.reltype,
                    value._target,
                    value.rId,
                )

        return self

    def move_slide(
        self,
        slide_number: int,
        to_position: int,
    ) -> Self:
        raise NotImplementedError("Moving slides is not implemented yet.")
        return self

    def reorder_slides(
        self,
        new_order: list[int],
    ) -> Self:
        raise NotImplementedError("Reordering slides is not implemented yet.")
        return self

    def insert_textbox(
        self,
        slide_number: Optional[int] = None,
        height: Optional[Length] = None,
        width: Optional[Length] = None,
        x_shift: Optional[Length] = None,
        y_shift: Optional[Length] = None,
        **kwargs,
    ) -> Shape:
        if height is None:
            if self.height is None:
                raise ValueError(
                    "Height must be specified if presentation height is not set."
                )
            height = Length.from_float(self.height * 0.9)
        if width is None:
            if self.width is None:
                raise ValueError(
                    "Width must be specified if presentation width is not set."
                )
            width = Length.from_float(self.width * 0.9)
        if x_shift is None:
            if self.width is None:
                raise ValueError(
                    "Width must be specified if presentation width is not set."
                )
            x_shift = Length.from_float(self.width * 0.05)
        if y_shift is None:
            if self.height is None:
                raise ValueError(
                    "Height must be specified if presentation height is not set."
                )
            y_shift = Length.from_float(self.height * 0.05)

        slide = self.slide(
            slide_number=slide_number
            if slide_number is not None
            else len(self.slides) - 1
        )

        textbox = slide.shapes.add_textbox(
            left=x_shift,
            top=y_shift,
            width=width,
            height=height,
        )

        return textbox

    def insert_text(
        self,
        text: str,
        textbox: Shape,
        bold: bool = False,
        italic: bool = False,
        underline: bool = False,
        strikethrough: bool = False,
        font_size: Optional[int] = None,
        font_family: Optional[str] = None,
        colour: Optional[Colour | str] = None,
        background_colour: Optional[Colour | str] = None,
        link: Optional[str] = None,
    ) -> Self:
        if colour is not None and not isinstance(colour, Colour):
            colour = Colour.parse(colour=colour)
        if background_colour is not None and not isinstance(background_colour, Colour):
            background_colour = Colour.parse(colour=background_colour)

        textbox.text_frame.text = text

        if font_size is not None:
            textbox.text_frame.paragraphs[0].font.size = Pt(points=font_size)
        if font_family is not None:
            textbox.text_frame.paragraphs[0].font.name = font_family
        if colour is not None:
            textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                r=colour.r,
                g=colour.g,
                b=colour.b,
            )
        if background_colour is not None:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(
                r=background_colour.r,
                g=background_colour.g,
                b=background_colour.b,
            )
        if bold:
            textbox.text_frame.paragraphs[0].font.bold = True
        if italic:
            textbox.text_frame.paragraphs[0].font.italic = True
        if underline:
            textbox.text_frame.paragraphs[0].font.underline = True
        if strikethrough:
            textbox.text_frame.paragraphs[0].font._element.attrib["strike"] = (
                "sngStrike"
            )

        if link is not None:
            raise NotImplementedError("Jyperlinks are not implemented yet.")
            from pptx.util import URI

            textbox.text_frame.paragraphs[0].hyperlink.address = URI(link)

        return self

    def save(
        self,
        file_path: Path | str,
    ) -> None:
        file_path = Path(file_path)

        return self.internal.save(
            file=str(file_path),
        )
