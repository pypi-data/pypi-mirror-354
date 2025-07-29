# ====== Code Summary ======
# This module defines `PPTXExtractor`, a subclass of `BaseExtractor`, designed to extract text
# from PowerPoint `.pptx` files using `python-pptx`. It also includes a method `describe_image`
# to send images to an OpenAI model for captioning (currently commented out).
# The class logs all steps of the extraction and handles both text and image content.

REQUIRED_LIBS = ["python-pptx", "openai"]

# ====== Third-party Library Imports ======
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import openai

    _NOT_INSTALLED = False
except ImportError:
    _NOT_INSTALLED = True

# ====== Internal Project Imports ======
from extractor_wrapper.base_extractor import BaseExtractor


class PPTXExtractor(BaseExtractor):
    """
    Extractor for PPTX (PowerPoint) files.

    Inherits from:
        BaseExtractor: Provides file validation and logging.

    This extractor pulls textual content from each slide and optionally describes embedded images
    using an OpenAI image description model.
    """

    def describe_image(self, image_bytes: bytes) -> str:
        """
        Send image data to an OpenAI image captioning model and return the generated description.

        Args:
            image_bytes (bytes): The image data to describe.

        Returns:
            str: Generated description or fallback error message if an issue occurs.
        """
        self.logger.info("Sending image to model for description.")
        try:
            # Placeholder call — replace with actual API endpoint and parameters
            response = openai.Image.create(
                file=image_bytes,
                model="image-alpha-001",  # Replace with your actual model identifier.
                prompt="Describe the image"
            )
            description = response.get("data", [{}])[0].get("caption", "No description provided")
        except Exception as e:
            self.logger.error(f"Error during image description: {e}")
            description = "Error describing image"
        return description

    def _ext_extract(self, file_path: str) -> str:
        """
        Extract text content from a PPTX file. Optionally supports image processing (currently commented).

        Args:
            file_path (str): Path to the .pptx file.

        Returns:
            str: Combined extracted text from all slides.

        Raises:
            Exception: If extraction or file processing fails.
        """
        self.logger.info(f"Extracting content from PPTX: {file_path}")
        extracted_text = []
        try:
            presentation = Presentation(file_path)
            for slide_index, slide in enumerate(presentation.slides):
                self.logger.info(f"Processing slide {slide_index + 1}")
                for shape_index, shape in enumerate(slide.shapes):
                    # Extract text if present in shape
                    if hasattr(shape, "text") and shape.text:
                        self.logger.debug(
                            f"Extracting text from shape {shape_index + 1} in slide {slide_index + 1}"
                        )
                        extracted_text.append(shape.text)

                    # Uncomment the following block to enable image extraction and description
                    # if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    #     self.logger.info(f"Found image in slide {slide_index + 1}, shape {shape_index + 1}")
                    #     try:
                    #         image_bytes = shape.image.blob
                    #         description = self.describe_image(image_bytes)
                    #         base_name = os.path.splitext(os.path.basename(file_path))[0]
                    #         image_file_name = f"{base_name}_slide{slide_index + 1}_img{shape_index + 1}.jpg.txt"
                    #         with open(image_file_name, "w", encoding="utf-8") as img_desc_file:
                    #             img_desc_file.write(description)
                    #         self.logger.info(f"Saved image description to {image_file_name}")
                    #     except Exception as img_e:
                    #         self.logger.error(
                    #             f"Error processing image in slide {slide_index + 1}, shape {shape_index + 1}: {img_e}"
                    #         )

            return "\n".join(extracted_text)

        except Exception as e:
            self.logger.error(f"Error extracting PPTX: {file_path}: {e}")
            raise e
